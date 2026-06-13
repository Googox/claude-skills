#!/usr/bin/env python3
"""
ARENA Executive Search — Businessplan 2026–2029
Apple-Style: Weiß / Schwarz / Grau — kein Farbenspiel
Aaron Arena | Start: 01.10.2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SW = Inches(13.33)
SH = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# DESIGN-SYSTEM — NUR SCHWARZ / GRAU / WEISS
# ─────────────────────────────────────────────────────────────────────────────

WHITE   = RGBColor(0xFF, 0xFF, 0xFF)   # Hintergrund
BK      = RGBColor(0x1D, 0x1D, 0x1F)  # Apple-Schwarz — Überschriften
DK      = RGBColor(0x3A, 0x3A, 0x3C)  # Dunkelgrau — fließtext
MD      = RGBColor(0x6E, 0x6E, 0x73)  # Mittelgrau — Labels
LT      = RGBColor(0xF5, 0xF5, 0xF7)  # Hellgrau — Kartenhintergrund
VLT     = RGBColor(0xFA, 0xFA, 0xFC)  # Sehr hell — alternate rows
DIV     = RGBColor(0xD2, 0xD2, 0xD7)  # Trennlinie
XDIV    = RGBColor(0xB0, 0xB0, 0xB5)  # Etwas dunklere Linie

F  = 'Calibri'
FL = 'Calibri Light'

_id = [0]
def _nid():
    _id[0] += 1
    return _id[0]
def _rst():
    _id[0] = 0

# ─────────────────────────────────────────────────────────────────────────────
# PRIMITIVE HELPER
# ─────────────────────────────────────────────────────────────────────────────

def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill; f.solid(); f.fore_color.rgb = WHITE
    return s

def rect(s, x, y, w, h, fill=None, lc=None, lw=0.5):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if lc:   sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, sz=16, clr=None, bold=False, italic=False,
        align=PP_ALIGN.LEFT, fn=F, wrap=True):
    if clr is None: clr = BK
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.name = fn; r.font.size = Pt(sz)
    r.font.color.rgb = clr; r.font.bold = bold; r.font.italic = italic
    return tb

def txb(s, rows, x, y, w, h, wrap=True):
    """rows = list of (text, sz, clr, bold, align, italic, space_before?)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    for i, row in enumerate(rows):
        text, sz, clr, bold, align, italic = row[:6]
        sp = row[6] if len(row) > 6 else 0
        p  = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if sp: p.space_before = Pt(sp)
        r  = p.add_run()
        r.text = text; r.font.name = F; r.font.size = Pt(sz)
        r.font.color.rgb = clr or BK
        r.font.bold = bold; r.font.italic = italic
    return tb

def slide_hdr(s, title, sub=None):
    """Standard-Folientitel mit dünner schwarzer Linie."""
    txt(s, title, 0.55, 0.32, 12.3, 0.58, sz=28, clr=BK, fn=FL)
    rect(s, 0.55, 0.94, 12.23, 0.022, fill=DIV)
    if sub:
        txt(s, sub, 0.55, 1.0, 12.0, 0.3, sz=12, clr=MD, italic=True)

def footer(s, n, total=16):
    rect(s, 0, 7.28, 13.33, 0.22, fill=LT)
    txt(s, "ARENA Executive Search  —  Businessplan 2026–2029", 0.45, 7.3, 9, 0.18,
        sz=8, clr=MD)
    txt(s, f"{n} / {total}", 12.55, 7.3, 0.7, 0.18, sz=8, clr=MD, align=PP_ALIGN.RIGHT)

def fade_tr(slide):
    ns  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ex  = slide.element.find(f'{{{ns}}}transition')
    if ex is not None: slide.element.remove(ex)
    slide.element.append(etree.fromstring(
        f'<p:transition xmlns:p="{ns}" spd="fast"><p:fade/></p:transition>'))

# ─────────────────────────────────────────────────────────────────────────────
# ANIMATION BUILDER
# ─────────────────────────────────────────────────────────────────────────────

def apply_anim(slide, specs):
    """specs: [(shape_id, 'fade'|'fly_up', click_grp, delay_ms, dur_ms)]"""
    if not specs: return
    _rst()
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    from collections import defaultdict
    grps = defaultdict(list)
    for sp_id, anim, grp, delay, dur in specs:
        grps[grp].append((sp_id, anim, delay, dur))

    def E(tag, **a):
        e = etree.Element(f'{{{ns}}}{tag}')
        for k, v in a.items(): e.set(k, str(v))
        return e

    timing = E('timing')
    tnLst  = E('tnLst'); timing.append(tnLst)
    par0   = E('par');   tnLst.append(par0)
    cTn0   = E('cTn', id=_nid(), dur='indefin', restart='whenNotActive', nodeType='tmRoot')
    par0.append(cTn0)
    ch0    = E('childTnLst'); cTn0.append(ch0)
    seq    = E('seq', concurrent='1', nextAc='seek'); ch0.append(seq)
    cTnS   = E('cTn', id=_nid(), dur='indefin', nodeType='mainSeq'); seq.append(cTnS)
    chS    = E('childTnLst'); cTnS.append(chS)

    for grp_idx in sorted(grps.keys()):
        items = grps[grp_idx]
        p_cg  = E('par'); chS.append(p_cg)
        cTn_cg = E('cTn', id=_nid(), fill='hold'); p_cg.append(cTn_cg)
        sc1 = E('stCondLst'); cTn_cg.append(sc1); sc1.append(E('cond', delay='indefin'))
        ch_cg = E('childTnLst'); cTn_cg.append(ch_cg)
        p_in  = E('par'); ch_cg.append(p_in)
        cTn_in = E('cTn', id=_nid(), fill='hold'); p_in.append(cTn_in)
        sc2 = E('stCondLst'); cTn_in.append(sc2); sc2.append(E('cond', delay='0'))
        ch_in = E('childTnLst'); cTn_in.append(ch_in)

        for i, (sp_id, anim, delay, dur) in enumerate(items):
            nt = 'clickEffect' if i == 0 else 'withEffect'
            pa = E('par'); ch_in.append(pa)
            pr = '10' if anim == 'fade' else '2'
            ps = '0'  if anim == 'fade' else '8'
            ca = E('cTn', id=_nid(), presetID=pr, presetClass='entr',
                   presetSubtype=ps, fill='hold', grpId=str(grp_idx-1), nodeType=nt)
            pa.append(ca)
            sc3 = E('stCondLst'); ca.append(sc3); sc3.append(E('cond', delay=str(delay)))
            ch_a = E('childTnLst'); ca.append(ch_a)
            # visibility
            se = E('set'); ch_a.append(se)
            cb = E('cBhvr'); se.append(cb)
            cb.append(E('cTn', id=_nid(), dur='1'))
            te = E('tgtEl'); cb.append(te); te.append(E('spTgt', spid=str(sp_id)))
            al = E('attrNameLst'); cb.append(al)
            an = E('attrName'); an.text = 'style.visibility'; al.append(an)
            to = E('to'); se.append(to); sv = E('strVal', val='visible'); to.append(sv)
            # effect
            flt = 'fade' if anim == 'fade' else 'wipe(up)'
            ae  = E('animEffect', transition='in', filter=flt); ch_a.append(ae)
            cb2 = E('cBhvr'); ae.append(cb2)
            cb2.append(E('cTn', id=_nid(), dur=str(dur)))
            te2 = E('tgtEl'); cb2.append(te2); te2.append(E('spTgt', spid=str(sp_id)))

    pv = E('prevCondLst'); seq.append(pv)
    cp = E('cond', evt='onPrevClick', delay='0'); pv.append(cp); cp.append(E('tn'))
    timing.append(E('bldLst'))

    ns_str = f'{{{ns}}}'
    ex = slide.element.find(f'{ns_str}timing')
    if ex is not None: slide.element.remove(ex)
    slide.element.append(timing)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER (weiß, typografisch, Apple-minimalistisch)
# ─────────────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = new_slide(prs)
    # Schmale schwarze Linie links — einzige Dekoration
    rect(s, 0.55, 1.6, 0.04, 4.3, fill=BK)

    t1 = txt(s, "ARENA", 0.9, 1.55, 11, 1.6, sz=96, clr=BK, fn=FL, bold=False)
    t2 = txt(s, "Executive Search", 0.9, 3.1, 11, 0.7, sz=34, clr=MD, fn=FL)
    ln = rect(s, 0.9, 3.9, 7.0, 0.025, fill=DIV)
    t3 = txt(s, "Die richtigen Führungspersönlichkeiten.\nZur richtigen Zeit.", 0.9, 4.05, 9, 0.7,
             sz=17, clr=DK, fn=FL)
    t4 = txt(s, "Businessplan  2026 – 2029  |  Aaron Arena  |  Oktober 2026",
             0.9, 6.75, 10, 0.3, sz=10, clr=MD)

    fade_tr(s)
    apply_anim(s, [
        (t1.shape_id, 'fade', 1,    0, 800),
        (t2.shape_id, 'fade', 1,  700, 600),
        (ln.shape_id, 'fade', 1, 1100, 400),
        (t3.shape_id, 'fade', 1, 1400, 600),
        (t4.shape_id, 'fade', 1, 2000, 500),
    ])
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — INHALTSVERZEICHNIS
# ─────────────────────────────────────────────────────────────────────────────

def slide_agenda(prs):
    s = new_slide(prs)
    footer(s, 2)
    slide_hdr(s, "Inhaltsverzeichnis")

    items = [
        "Die Marktchance",
        "Der Gründer — Aaron Arena",
        "Marktanalyse DACH",
        "Geschäftsmodell & Einnahmen",
        "Leistungsportfolio",
        "Wettbewerbsvorteile (USP)",
        "Partnerschaftsmodell — Nachfolgelösung",
        "Wachstumsstrategie",
        "Finanzplanung 2026 – 2029",
        "Investitionsbedarf & Mittelverwendung",
        "Meilensteine & Roadmap",
        "Fazit & Nächste Schritte",
    ]

    for i, label in enumerate(items):
        col = i // 6; row = i % 6
        x = 0.55 + col * 6.5; y = 1.45 + row * 0.88
        num = f"{i+1:02d}"
        rect(s, x, y + 0.04, 0.4, 0.38, fill=BK)
        txt(s, num, x, y + 0.04, 0.4, 0.38, sz=11, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, label, x + 0.52, y + 0.07, 5.7, 0.35, sz=13, clr=BK)
        rect(s, x + 0.52, y + 0.48, 5.7, 0.012, fill=DIV)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — DIE MARKTCHANCE (grosse Zahlen, reines Weiß)
# ─────────────────────────────────────────────────────────────────────────────

def slide_chance(prs):
    s = new_slide(prs)
    footer(s, 3)
    slide_hdr(s, "Die Marktchance",
              "Executive Search in Deutschland — strukturelles Wachstum in einem hochmargigen Markt")

    # Drei grosse KPI-Blöcke
    kpis = [
        ("€ 2,8 Mrd.", "Executive-Search-Markt\nDeutschland 2026"),
        ("+ 11 %",      "Jährliches Wachstum\n(CAGR 2023–2028)"),
        ("67 %",        "C-Level-Positionen\nwerden diskret besetzt"),
    ]
    for i, (val, lbl) in enumerate(kpis):
        cx = 0.55 + i * 4.25
        rect(s, cx, 1.55, 4.0, 2.2, fill=LT)
        txt(s, val, cx + 0.2, 1.65, 3.6, 1.05, sz=48, clr=BK, fn=FL, bold=False,
            align=PP_ALIGN.CENTER)
        rect(s, cx + 0.3, 2.68, 3.4, 0.018, fill=DIV)
        txt(s, lbl, cx + 0.2, 2.73, 3.6, 0.75, sz=12, clr=MD,
            align=PP_ALIGN.CENTER, wrap=True)

    # Vier Treiber darunter
    drivers = [
        ("Digitale Transformation",
         "70 % der Unternehmen suchen CDO/CTO — Nachfrage +34 % YoY"),
        ("Nachfolgeplanung",
         "580.000 KMU-Übergaben bis 2030 — Strukturwandel als Dauertreiber"),
        ("Regulierung & ESG",
         "Compliance-Druck erhöht Nachfrage nach spezialisierten Risk-Executives"),
        ("Internationalisierung",
         "Cross-border Searches DACH +28 % — Grenzübergreifender Talent-Pool"),
    ]
    for i, (title, desc) in enumerate(drivers):
        cx = 0.55 + i * 3.2
        rect(s, cx, 4.05, 3.05, 1.55, fill=WHITE, lc=DIV, lw=0.5)
        rect(s, cx, 4.05, 3.05, 0.04, fill=BK)
        txt(s, title, cx + 0.12, 4.15, 2.8, 0.38, sz=11, clr=BK, bold=True)
        txt(s, desc,  cx + 0.12, 4.58, 2.8, 0.9,  sz=10, clr=MD, wrap=True)

    txt(s, "Quellen: BDU Executive Search Report 2026  |  Statista 2026  |  Kienbaum Executive Panel",
        0.55, 6.82, 12, 0.25, sz=8, clr=DIV, italic=True)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — DER GRÜNDER
# ─────────────────────────────────────────────────────────────────────────────

def slide_founder(prs):
    s = new_slide(prs)
    footer(s, 4)
    slide_hdr(s, "Der Gründer — Aaron Arena",
              "30+ Jahre Erfahrung in Executive Recruitment | Gründung 01.10.2026")

    # Linke Spalte — grosse Zahl + Highlights
    rect(s, 0.55, 1.38, 4.0, 5.65, fill=LT)
    txt(s, "30+", 0.7, 1.5, 3.7, 1.1, sz=72, clr=BK, fn=FL, align=PP_ALIGN.CENTER)
    txt(s, "Jahre im Executive Search", 0.7, 2.55, 3.7, 0.35, sz=12, clr=MD,
        align=PP_ALIGN.CENTER)
    rect(s, 1.1, 3.0, 2.9, 0.02, fill=DIV)

    highlights = [
        "200+  erfolgreiche C-Level-Placements",
        "1.500+  Entscheider-Kontakte DACH",
        "Sektorfokus: Finance · Tech · Industrie",
        "Ausbildung: BWL, Executive Development",
        "Sprachen: Deutsch · Englisch · Italienisch",
    ]
    for i, h in enumerate(highlights):
        txt(s, h, 0.7, 3.15 + i * 0.5, 3.65, 0.42, sz=11, clr=DK)

    # Rechte Spalte — Karriere-Timeline
    timeline = [
        ("2026 →", "Gründer", "ARENA Executive Search — München",    True),
        ("2019–2026", "Senior Partner",   "Boutique Executive Search — DACH", False),
        ("2012–2019", "Director Search",  "Internationale Personalberatung — Frankfurt", False),
        ("2005–2012", "Senior Consultant","Finanz- & Technologiesektor — DACH", False),
        ("1995–2005", "Karrierebeginn",   "Unternehmensberatung & Personalwesen", False),
    ]

    for i, (year, role, desc, current) in enumerate(timeline):
        yp = 1.42 + i * 1.02
        lc = BK if current else DIV
        rect(s, 4.9, yp + 0.04, 0.04, 0.85, fill=lc)
        rect(s, 4.77, yp + 0.1, 0.2, 0.2, fill=BK if current else DIV)
        txt(s, year, 5.15, yp,        3.5, 0.3,  sz=9,  clr=BK if current else MD, bold=current)
        txt(s, role, 5.15, yp + 0.28, 7.6, 0.32, sz=13, clr=BK, bold=True)
        txt(s, desc, 5.15, yp + 0.58, 7.6, 0.34, sz=11, clr=MD)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — MARKTANALYSE DACH
# ─────────────────────────────────────────────────────────────────────────────

def slide_markt(prs):
    s = new_slide(prs)
    footer(s, 5)
    slide_hdr(s, "Marktanalyse — Executive Search DACH",
              "Kontinuierliches Wachstum in einem strukturell robusten Markt")

    # Balkendiagramm links
    rect(s, 0.55, 1.38, 6.15, 5.35, fill=LT)
    txt(s, "Marktvolumen Executive Search DACH (Mrd. €)", 0.7, 1.5, 5.8, 0.32,
        sz=11, clr=DK, bold=True)

    bars = [
        ("2022", 2.0), ("2023", 2.2), ("2024", 2.5),
        ("2025e", 2.65), ("2026e", 2.8), ("2027p", 3.1),
    ]
    max_val = 3.4; max_h = 2.8; bw = 0.62
    for i, (yr, val) in enumerate(bars):
        bh = (val / max_val) * max_h
        bx = 0.75 + i * 0.88; by = 5.6 - bh
        fill = BK if i >= 3 else DK
        rect(s, bx, by, bw, bh, fill=fill)
        txt(s, f"{val:.1f}", bx - 0.02, by - 0.3, 0.7, 0.25,
            sz=8, clr=BK if i >= 3 else MD, align=PP_ALIGN.CENTER, bold=(i >= 3))
        txt(s, yr, bx - 0.01, 5.65, 0.68, 0.22,
            sz=9, clr=DK, align=PP_ALIGN.CENTER)

    txt(s, "e = geschätzt  |  p = Prognose  |  Quelle: BDU, Statista 2026",
        0.7, 6.6, 5.8, 0.2, sz=8, clr=MD, italic=True)

    # Rechts: Marktsegmente
    txt(s, "Marktsegmente nach Funktion", 7.0, 1.38, 5.9, 0.32, sz=11, clr=DK, bold=True)
    rect(s, 7.0, 1.73, 5.9, 0.022, fill=DIV)

    segs = [
        ("CEO / Geschäftsführung",    "38 %"),
        ("CFO / Finance Leadership",  "22 %"),
        ("CTO / Digital / CDO",       "18 %"),
        ("COO / Operations",          "12 %"),
        ("Sonstige C-Suite",          "10 %"),
    ]
    for i, (label, pct) in enumerate(segs):
        ry = 1.88 + i * 0.72
        bar_w = float(pct.rstrip(" %")) / 100 * 4.8
        rect(s, 7.0, ry, bar_w, 0.32, fill=BK if i == 0 else (DK if i < 3 else MD))
        txt(s, label, 7.08, ry + 0.06, 3.5, 0.24, sz=10, clr=WHITE)
        txt(s, pct,   11.7, ry + 0.06, 0.5, 0.24, sz=10, clr=DK, bold=True)

    # Rechts unten: Insights
    insights = [
        "Nachfolgemandate: +31 % gegenüber Vorjahr",
        "Frauenanteil C-Level: Regulierung als Treiber",
        "Ø Suchdauer Top-Mandat: 12–16 Wochen",
        "Boutiques gewinnen Marktanteile von Big 4",
    ]
    for i, ins in enumerate(insights):
        rect(s, 7.0, 5.5 + i * 0.4, 5.9, 0.36, fill=LT if i % 2 == 0 else WHITE)
        txt(s, f"→  {ins}", 7.1, 5.54 + i * 0.4, 5.7, 0.28, sz=10, clr=DK)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — GESCHÄFTSMODELL
# ─────────────────────────────────────────────────────────────────────────────

def slide_modell(prs):
    s = new_slide(prs)
    footer(s, 6)
    slide_hdr(s, "Geschäftsmodell & Einnahmen",
              "Drei Säulen — Retained Search als Kern")

    cols = [
        ("Retained\nSearch", "Kernleistung",
         ["Exklusives Mandat", "Retainer 1/3 upfront + Milestone + Success",
          "Gebühr: 28–33 % Jahresgehalt", "Garantiezeitraum: 12 Monate",
          "Ø Laufzeit: 10–14 Wochen"],
         "60 %", BK),
        ("Executive\nInterim", "Ergänzend",
         ["Sofortbesetzung bei Führungslücken", "3–12 Monate Laufzeit",
          "Tagesatz: €1.200 – 2.200", "Kein Exklusiv-Mandat nötig",
          "Schneller ROI für den Mandanten"],
         "25 %", DK),
        ("Advisory &\nAssessment", "Ergänzend",
         ["C-Suite Potenzialanalysen", "Nachfolgeplanung intern",
          "Executive Coaching", "Workshop-Formate",
          "Festpreis: €8.000 – 18.000"],
         "15 %", MD),
    ]

    for i, (title, badge, bullets, share, hclr) in enumerate(cols):
        cx = 0.55 + i * 4.25
        rect(s, cx, 1.38, 4.0, 5.35, fill=LT if i % 2 == 0 else WHITE, lc=DIV, lw=0.5)
        rect(s, cx, 1.38, 4.0, 0.5,  fill=hclr)
        txt(s, title, cx + 0.15, 1.4, 3.7, 0.48, sz=18, clr=WHITE, bold=True,
            fn=FL, wrap=True, align=PP_ALIGN.CENTER)
        rect(s, cx + 0.8, 1.93, 2.4, 0.3, fill=DIV)
        txt(s, badge, cx + 0.8, 1.93, 2.4, 0.3, sz=9, clr=DK, bold=True,
            align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            txt(s, f"· {b}", cx + 0.18, 2.36 + j * 0.54, 3.65, 0.48,
                sz=11, clr=DK, wrap=True)
        rect(s, cx, 6.43, 4.0, 0.3, fill=hclr)
        txt(s, f"Umsatzanteil: {share}", cx + 0.1, 6.45, 3.8, 0.26,
            sz=10, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — LEISTUNGSPORTFOLIO
# ─────────────────────────────────────────────────────────────────────────────

def slide_leistungen(prs):
    s = new_slide(prs)
    footer(s, 7)
    slide_hdr(s, "Leistungsportfolio", "Was wir für unsere Mandanten leisten")

    services = [
        ("C-Suite Search",
         "Vorstand · Geschäftsführung · C-Level",
         "Vollumfänglicher, diskreter Suchprozess für höchste Führungsebenen. "
         "Longlist → Qualifizierung → Strukturiertes Assessment → Onboarding. "
         "Erfolgsquote > 95 %. 12 Monate Garantiezeitraum."),
        ("Board Advisory",
         "Aufsichtsrat · Beirat · Advisory Board",
         "Rekrutierung und Beratung für Kontrollgremien. Governance, Diversität "
         "und Kompetenzprofil-Analyse. Wachsend regulatorische Anforderungen als Treiber."),
        ("Executive Interim",
         "Überbrückung & Transformation",
         "Kurzfristig verfügbare Top-Führungskräfte für Restrukturierungen, "
         "M&A-Projekte und Vakanzüberbrückungen. Netzwerk von 120+ geprüften Interim-Executives."),
        ("Succession Planning",
         "Nachfolge & Talentpipeline",
         "Systematische Nachfolgeplanung für Inhaber, Familienunternehmen und "
         "PE-Portfolios. Begleitung über 12–36 Monate Transitionszeitraum."),
    ]

    for i, (title, sub, desc) in enumerate(services):
        col = i % 2; row = i // 2
        cx = 0.55 + col * 6.4; cy = 1.38 + row * 2.72
        rect(s, cx, cy, 6.1, 2.55, fill=LT if i % 2 == 0 else WHITE, lc=DIV, lw=0.5)
        rect(s, cx, cy, 0.05, 2.55, fill=BK)
        txt(s, title, cx + 0.2, cy + 0.1, 5.7, 0.38, sz=15, clr=BK, bold=True)
        txt(s, sub,   cx + 0.2, cy + 0.5, 5.7, 0.28, sz=10, clr=MD, italic=True)
        rect(s, cx + 0.2, cy + 0.82, 5.6, 0.018, fill=DIV)
        txt(s, desc,  cx + 0.2, cy + 0.9, 5.75, 1.4, sz=11, clr=DK, wrap=True)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — WETTBEWERBSVORTEILE (weiß, schwarz — kein Farbhintergrund mehr)
# ─────────────────────────────────────────────────────────────────────────────

def slide_usp(prs):
    s = new_slide(prs)
    footer(s, 8)
    slide_hdr(s, "Warum ARENA?",
              "Vier strategische Differenzierungsmerkmale gegenüber dem Wettbewerb")

    usps = [
        ("01", "Tiefes persönliches\nBranchennetzwerk",
         "30+ Jahre aufgebautes Netzwerk von 1.500+ C-Level-Kontakten in Finance, Technology "
         "und Industrie im DACH-Raum. Direktansprache — kein anonymes Datenbanksurf."),
        ("02", "Boutique-Qualität\nbei jeder Suche",
         "Kein Volumenziel. Kein Staffing. Jede Suche wird vom Gründer persönlich geführt. "
         "Mandantenzahl bewusst limitiert für maximale Qualität und volle Aufmerksamkeit."),
        ("03", "Nachfolge-Expertise\nals Alleinstellungsmerkmal",
         "Spezialkompetenz in Unternehmensnachfolge und Senior-Transition. "
         "Wachsendes Marktsegment: 580.000 KMU-Nachfolgen bis 2030 in Deutschland."),
        ("04", "Diskrete, vertrauens-\nbasierte Arbeitsweise",
         "Vertraulichkeit als Kernprinzip. Keine Stellenanzeigen, keine öffentliche Sichtbarkeit. "
         "Preferred Partner zahlreicher Familienunternehmen und Private-Equity-Häuser."),
    ]

    for i, (num, title, desc) in enumerate(usps):
        cx = 0.55 + (i % 2) * 6.4
        cy = 1.38 + (i // 2) * 2.75
        rect(s, cx, cy, 6.1, 2.55, fill=LT if i % 2 == 0 else WHITE, lc=DIV, lw=0.5)
        # Schwarze Nummer
        rect(s, cx + 0.15, cy + 0.12, 0.55, 0.55, fill=BK)
        txt(s, num, cx + 0.15, cy + 0.12, 0.55, 0.55, sz=14, clr=WHITE, bold=True,
            align=PP_ALIGN.CENTER)
        txt(s, title, cx + 0.85, cy + 0.1, 5.1, 0.7, sz=14, clr=BK, bold=True,
            fn=FL, wrap=True)
        rect(s, cx + 0.85, cy + 0.88, 4.9, 0.018, fill=DIV)
        txt(s, desc, cx + 0.2, cy + 0.98, 5.7, 1.35, sz=11, clr=DK, wrap=True)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — PARTNERSCHAFTSMODELL
# ─────────────────────────────────────────────────────────────────────────────

def slide_partner(prs):
    s = new_slide(prs)
    footer(s, 9)
    slide_hdr(s, "Partnerschaftsmodell — Nachfolgelösung",
              "Strukturierte Praxisübernahme über 4 Jahre — Win-Win für beide Seiten")

    # Erklärungstext links
    rect(s, 0.55, 1.38, 5.7, 5.35, fill=LT)
    txt(s, "Das Modell", 0.7, 1.48, 5.4, 0.32, sz=13, clr=BK, bold=True)
    rect(s, 0.7, 1.82, 5.3, 0.02, fill=DIV)
    desc = (
        "Ein erfahrener Headhunter (Gründer, 60+) sucht einen qualifizierten "
        "Nachfolger, der sein Lebenswerk weiterführt und über 4 Jahre die Mandanten, "
        "das Know-how und den Markennamen übernimmt.\n\n"
        "Aaron Arena als designierter Nachfolger bringt mit:\n"
        "· Operatives Executive-Search-Know-how (30+ Jahre)\n"
        "· Eigenständiges DACH-Netzwerk (1.500+ Kontakte)\n"
        "· Finanzielle Kapazität & Bankfinanzierung\n"
        "· Wachstumskonzept für die bestehende Praxis\n\n"
        "Der Senior-Partner erhält:\n"
        "· Strukturierten Ausstieg über 48 Monate\n"
        "· Faire Beteiligungsvergütung\n"
        "· Sicherung seines Lebenswerks"
    )
    txt(s, desc, 0.7, 1.95, 5.3, 4.6, sz=11, clr=DK, wrap=True)

    # 4 Phasen rechts
    phases = [
        ("Phase 1  |  2026–2027",  "Einstieg & Lernen",
         "Shadowing, gemeinsame Mandate, dual-brand"),
        ("Phase 2  |  2027–2028",  "Operative Übernahme",
         "Aaron Arena führt 60 % der Mandate eigenständig"),
        ("Phase 3  |  2028–2029",  "Mehrheitskontrolle",
         "Übernahme der Mandantschaft & Beteiligung"),
        ("Phase 4  |  2030",       "Vollübernahme",
         "100 % ARENA — Senior-Partner in Advisory-Rolle"),
    ]
    fill_seq = [BK, DK, MD, DK]
    for i, (ph, title, desc2) in enumerate(phases):
        cy = 1.38 + i * 1.22
        rect(s, 6.55, cy, 6.5, 1.12, fill=LT if i % 2 == 0 else WHITE, lc=DIV, lw=0.5)
        rect(s, 6.55, cy, 6.5, 0.36, fill=fill_seq[i])
        txt(s, ph,    6.68, cy + 0.06, 6.2, 0.26, sz=10, clr=WHITE, bold=True)
        txt(s, title, 6.68, cy + 0.44, 6.1, 0.3,  sz=13, clr=BK, bold=True)
        txt(s, desc2, 6.68, cy + 0.76, 6.1, 0.28, sz=11, clr=MD)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — WACHSTUMSSTRATEGIE
# ─────────────────────────────────────────────────────────────────────────────

def slide_wachstum(prs):
    s = new_slide(prs)
    footer(s, 10)
    slide_hdr(s, "Wachstumsstrategie", "Fokussiertes Wachstum in drei Horizonten")

    horizons = [
        ("H1  —  Aufbau",   "Okt 2026 – Sep 2027",
         ["5–8 Retained Searches abschließen", "3 Anker-Mandanten gewinnen",
          "Website, Brand, Reputation aufbauen", "Partnerschaft formalisieren",
          "Bestandsnetzwerk vollständig aktivieren"],
         "€ 240K", BK),
        ("H2  —  Wachstum", "2028",
         ["12–15 Mandate p.a.", "Spezialisierung Finance / Tech vertieft",
          "2. Consultant hinzunehmen", "Nachfolge-Boutique-Marke etabliert",
          "Erste Interim-Mandate produktiv"],
         "€ 520K", DK),
        ("H3  —  Skalierung","2029",
         ["18–22 Mandate p.a.", "Vollübernahme Partnerpraxis",
          "Board-Advisory-Segment ausgebaut", "DACH-weite Mandantenbasis",
          "Team: 3–4 Berater"],
         "€ 980K", MD),
    ]

    for i, (title, period, bullets, kpi, hclr) in enumerate(horizons):
        cx = 0.55 + i * 4.25
        rect(s, cx, 1.38, 4.0, 5.35, fill=LT if i%2==0 else WHITE, lc=DIV, lw=0.5)
        rect(s, cx, 1.38, 4.0, 0.55, fill=hclr)
        txt(s, title,  cx + 0.15, 1.4,  3.7, 0.3,  sz=14, clr=WHITE, bold=True)
        txt(s, period, cx + 0.15, 1.72, 3.7, 0.22, sz=10, clr=RGBColor(0xC8,0xC8,0xCC))
        for j, b in enumerate(bullets):
            txt(s, f"· {b}", cx + 0.18, 2.1 + j * 0.6, 3.65, 0.52,
                sz=11, clr=DK, wrap=True)
        rect(s, cx, 6.45, 4.0, 0.28, fill=hclr)
        txt(s, f"Ziel-Umsatz:  {kpi}", cx + 0.15, 6.47, 3.7, 0.24,
            sz=11, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)

    rect(s, 0.55, 6.82, 12.4, 0.28, fill=LT)
    txt(s, "Akquise-Strategie:  Bestandsnetzwerk  →  Empfehlungen  →  "
           "Verbandspräsenz  →  LinkedIn Thought-Leadership  →  Persönliche Direktansprache",
        0.7, 6.85, 12.0, 0.22, sz=9, clr=MD)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — FINANZPLANUNG
# ─────────────────────────────────────────────────────────────────────────────

def slide_finanzen(prs):
    s = new_slide(prs)
    footer(s, 11)
    slide_hdr(s, "Finanzplanung 2026 – 2029",
              "Konservative Prognose  |  Alle Werte in EUR  |  vor Steuern")

    years  = ["", "Q4 2026", "GJ 2027", "GJ 2028", "GJ 2029"]
    col_xs = [0.55, 4.3, 6.42, 8.54, 10.66]
    col_ws = [3.7, 2.08, 2.08, 2.08, 2.58]

    # Header
    for i, (label, cx, cw) in enumerate(zip(years, col_xs, col_ws)):
        rect(s, cx, 1.38, cw, 0.42, fill=BK if i > 0 else LT)
        txt(s, label, cx + 0.05, 1.38, cw - 0.05, 0.42, sz=12,
            clr=WHITE if i > 0 else MD, bold=True,
            align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    rows = [
        ("Placements (Anzahl)",      ["3",       "8",       "15",       "22"],       False, False),
        ("Ø Honorar (€)",            ["50.000",  "55.000",  "58.000",   "60.000"],   False, False),
        ("Retained-Search-Umsatz",   ["150.000", "440.000", "870.000",  "1.320.000"],True,  False),
        ("+ Interim / Advisory",     ["90.000",  "80.000",  "110.000",  "160.000"],  False, False),
        ("= Gesamtumsatz",           ["240.000", "520.000", "980.000",  "1.480.000"],True,  True),
        ("Betriebskosten",           ["130.000", "170.000", "220.000",  "310.000"],  False, False),
        ("EBIT",                     ["110.000", "350.000", "760.000",  "1.170.000"],True,  True),
        ("EBIT-Marge",               ["46 %",    "67 %",    "78 %",     "79 %"],     False, False),
    ]

    sep_rows = {4, 6}
    for ri, (label, vals, bold_row, highlight) in enumerate(rows):
        ry = 1.85 + ri * 0.56
        bg = LT if highlight else (VLT if ri % 2 == 0 else WHITE)
        rect(s, 0.55, ry, 12.69, 0.52, fill=bg)
        if ri in sep_rows:
            rect(s, 0.55, ry, 12.69, 0.025, fill=BK)
        txt(s, label, 0.65, ry + 0.08, 3.55, 0.38, sz=11,
            clr=BK if bold_row else DK, bold=bold_row)
        for ci, val in enumerate(vals):
            cx2 = col_xs[ci+1] + 0.05
            cw2 = col_ws[ci+1] - 0.1
            txt(s, val, cx2, ry + 0.08, cw2, 0.38, sz=11,
                clr=BK if bold_row else DK, bold=bold_row,
                align=PP_ALIGN.RIGHT)

    txt(s, "* Retained-Search-Umsatz = Placements × Ø-Honorar  |  Interim & Advisory separat",
        0.55, 6.82, 12.4, 0.25, sz=8, clr=MD, italic=True)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — INVESTITIONSBEDARF
# ─────────────────────────────────────────────────────────────────────────────

def slide_invest(prs):
    s = new_slide(prs)
    footer(s, 12)
    slide_hdr(s, "Investitionsbedarf & Mittelverwendung",
              "Finanzierungsanfrage an die Bank  |  Gesamtbedarf: € 150.000")

    # Links: Aufstellung
    rect(s, 0.55, 1.38, 6.2, 5.35, fill=LT)
    txt(s, "Mittelverwendung", 0.7, 1.48, 5.9, 0.32, sz=13, clr=BK, bold=True)
    rect(s, 0.7, 1.83, 5.9, 0.02, fill=DIV)

    items = [
        ("Büro & Coworking  (12 Monate)",  "18.000 €"),
        ("IT, CRM & Research-Tools",        "12.000 €"),
        ("Marke, Website, Marketing",       "15.000 €"),
        ("Verbände & Netzwerk-Events",       "8.000 €"),
        ("Versicherungen & Recht",           "7.000 €"),
        ("Betriebskapital (lfd. Kosten)",   "60.000 €"),
        ("Gründungsreserve (3 Monate)",     "30.000 €"),
    ]
    for i, (label, val) in enumerate(items):
        ry = 1.98 + i * 0.58
        rect(s, 0.7, ry, 5.9, 0.54, fill=WHITE if i%2==0 else LT)
        txt(s, label, 0.82, ry + 0.1, 4.1, 0.34, sz=11, clr=DK)
        txt(s, val,   4.8,  ry + 0.1, 1.65, 0.34, sz=11, clr=BK, bold=True,
            align=PP_ALIGN.RIGHT)
    rect(s, 0.7, 6.04, 5.9, 0.04, fill=BK)
    txt(s, "GESAMT", 0.82, 6.13, 2.0, 0.3, sz=12, clr=BK, bold=True)
    txt(s, "150.000 €", 4.8, 6.13, 1.65, 0.3, sz=12, clr=BK, bold=True,
        align=PP_ALIGN.RIGHT)

    # Rechts: Konditionen
    rect(s, 7.05, 1.38, 6.0, 5.35, fill=WHITE, lc=DIV, lw=0.5)
    txt(s, "Vorgeschlagene Konditionen", 7.2, 1.48, 5.7, 0.32, sz=13, clr=BK, bold=True)
    rect(s, 7.2, 1.83, 5.7, 0.02, fill=DIV)

    conds = [
        ("Kreditbetrag",        "€ 150.000"),
        ("Laufzeit",            "60 Monate (5 Jahre)"),
        ("Tilgungsfrei",        "12 Monate"),
        ("Rückzahlung ab",      "Oktober 2027"),
        ("Monatliche Rate",     "ca. € 2.900 ab M13"),
        ("Sicherheit",          "Abtretung Forderungen +\nPersonalbürgschaft"),
        ("Break-Even",          "Monat 8  (Juni 2027)"),
    ]
    for i, (label, val) in enumerate(conds):
        ry = 1.98 + i * 0.58
        rect(s, 7.2, ry, 5.7, 0.54, fill=LT if i%2==0 else WHITE)
        txt(s, label, 7.32, ry + 0.1, 2.4, 0.34, sz=10, clr=MD)
        txt(s, val,   9.75, ry + 0.08, 3.0, 0.4,  sz=11, clr=BK, bold=True, wrap=True)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — RISIKEN & MITIGATION
# ─────────────────────────────────────────────────────────────────────────────

def slide_risiken(prs):
    s = new_slide(prs)
    footer(s, 13)
    slide_hdr(s, "Risiken & Mitigationsmaßnahmen",
              "Transparente Betrachtung — strukturierte Gegenmaßnahmen")

    risks = [
        ("Lange Anlaufphase",
         "Hoch", "Mittel",
         "Bestandsnetzwerk ab Tag 1 aktiv; qualifizierte Pipeline aus "
         "10+ Interessenten bereits vorhanden vor Start."),
        ("Konjunktureller Abschwung",
         "Mittel", "Hoch",
         "Executive Search bleibt in Krisen kritisch; Krisen erhöhen "
         "Nachfolgebedarf; Diversifizierung über Sektoren abfedert."),
        ("Ausfall Partnerkooperation",
         "Niedrig", "Hoch",
         "Notarielle Regelung im Partnerschaftsvertrag; "
         "ARENA funktioniert vollständig als Solo-Praxis ab Tag 1."),
        ("Intensiver Wettbewerb",
         "Mittel", "Niedrig",
         "Differenzierung über Persönlichkeit, Netzwerk & Boutique-Service. "
         "Kunden wechseln zu Boutiques aus Qualitätsgründen, nicht wegen Preis."),
    ]

    labels = {"Hoch": BK, "Mittel": DK, "Niedrig": MD}
    for i, (risk, warsch, impact, mitigation) in enumerate(risks):
        col = i % 2; row = i // 2
        cx = 0.55 + col * 6.4; cy = 1.38 + row * 2.78
        rect(s, cx, cy, 6.1, 2.6, fill=LT if i%2==0 else WHITE, lc=DIV, lw=0.5)
        txt(s, risk, cx + 0.2, cy + 0.1, 5.7, 0.38, sz=15, clr=BK, bold=True)
        # Wahrscheinlichkeit / Impact Badges
        for j, (key, val) in enumerate([("Wahrsch.", warsch), ("Impact", impact)]):
            bx = cx + 0.2 + j * 2.0
            rect(s, bx, cy + 0.58, 1.8, 0.28, fill=labels[val])
            txt(s, f"{key}: {val}", bx + 0.05, cy + 0.58, 1.7, 0.28,
                sz=9, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
        rect(s, cx + 0.2, cy + 0.96, 5.7, 0.018, fill=DIV)
        txt(s, mitigation, cx + 0.2, cy + 1.04, 5.7, 1.38, sz=11, clr=DK, wrap=True)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — MEILENSTEINE & ROADMAP
# ─────────────────────────────────────────────────────────────────────────────

def slide_roadmap(prs):
    s = new_slide(prs)
    footer(s, 14)
    slide_hdr(s, "Meilensteine & Roadmap", "Klare Etappenziele — Quartal für Quartal")

    milestones = [
        ("Q4\n2026", "Start",
         ["Gründung & Anmeldung", "Website live", "Erste 2 Mandate aktiv"]),
        ("Q1\n2027", "Aufbau",
         ["3 Retainer-Mandate laufend", "Partnervertrag unterzeichnet", "Break-Even in Sicht"]),
        ("Q2\n2027", "Wachsen",
         ["5. Placement", "Erstes Interim-Mandat", "Verbandsmitgliedschaft aktiv"]),
        ("Q3\n2027", "Profil",
         ["8. Placement", "10 aktive Mandanten", "Break-Even erreicht"]),
        ("GJ\n2028", "Skalieren",
         ["15 Placements", "2. Consultant", "Nachfolgemandate aktiv"]),
        ("GJ\n2029", "Führen",
         ["22 Placements", "Vollübernahme Praxis", "EBIT > €1,1 Mio."]),
    ]

    fill_seq = [BK, BK, DK, DK, MD, MD]
    rect(s, 0.55, 3.9, 12.33, 0.04, fill=XDIV)

    for i, (period, stage, bullets) in enumerate(milestones):
        cx  = 0.55 + i * 2.1
        dot_x = cx + 0.83
        rect(s, dot_x, 3.77, 0.26, 0.26, fill=fill_seq[i])

        if i % 2 == 0:           # Oben
            cy = 1.38
            rect(s, cx, cy, 2.0, 2.35, fill=LT, lc=DIV, lw=0.5)
            rect(s, cx, cy, 2.0, 0.42, fill=fill_seq[i])
            txt(s, period, cx + 0.08, cy + 0.05, 0.85, 0.38,
                sz=11, clr=WHITE, bold=True, wrap=True)
            txt(s, stage, cx + 1.0, cy + 0.1, 0.85, 0.3,
                sz=10, clr=WHITE, bold=True)
            for j, b in enumerate(bullets):
                txt(s, f"· {b}", cx + 0.1, cy + 0.52 + j * 0.56, 1.82, 0.48,
                    sz=9, clr=DK, wrap=True)
        else:                     # Unten
            cy = 4.25
            rect(s, cx, cy, 2.0, 2.35, fill=LT, lc=DIV, lw=0.5)
            rect(s, cx, cy, 2.0, 0.42, fill=fill_seq[i])
            txt(s, period, cx + 0.08, cy + 0.05, 0.85, 0.38,
                sz=11, clr=WHITE, bold=True, wrap=True)
            txt(s, stage, cx + 1.0, cy + 0.1, 0.85, 0.3,
                sz=10, clr=WHITE, bold=True)
            for j, b in enumerate(bullets):
                txt(s, f"· {b}", cx + 0.1, cy + 0.52 + j * 0.56, 1.82, 0.48,
                    sz=9, clr=DK, wrap=True)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — ZIELBRANCHEN
# ─────────────────────────────────────────────────────────────────────────────

def slide_branchen(prs):
    s = new_slide(prs)
    footer(s, 15)
    slide_hdr(s, "Zielbranchen & Zielgruppen", "Fokussierter Sektormix — Tiefe vor Breite")

    sectors = [
        ("Financial\nServices",    "CFO · CRO · CEO",
         "Private Banking · PE · Versicherung · Asset Management"),
        ("Technology &\nDigital",  "CTO · CDO · CISO",
         "Software · FinTech · AI/ML · SaaS · Deep Tech"),
        ("Industrie &\nMittelstand","CEO · COO · CSO",
         "Familienunternehmen · Hidden Champions · Automotive"),
        ("Nachfolge & M&A",        "Übergabe · Integration",
         "Nachfolgeplanung · Post-Merger-Führung · PE-Portfolio"),
    ]

    fill_seq = [BK, DK, MD, BK]
    for i, (sector, roles, desc) in enumerate(sectors):
        cx = 0.55 + (i % 2) * 6.4
        cy = 1.38 + (i // 2) * 2.85
        rect(s, cx, cy, 6.1, 2.65, fill=LT if i%2==0 else WHITE, lc=DIV, lw=0.5)
        rect(s, cx, cy, 6.1, 0.52, fill=fill_seq[i])
        txt(s, sector, cx + 0.2, cy + 0.07, 5.7, 0.44, sz=17, clr=WHITE, bold=True,
            fn=FL, wrap=True)
        txt(s, "Zielpositionen", cx + 0.2, cy + 0.65, 2.0, 0.28,
            sz=9, clr=MD, italic=True)
        txt(s, roles, cx + 0.2, cy + 0.93, 5.7, 0.32, sz=12, clr=BK, bold=True)
        rect(s, cx + 0.2, cy + 1.27, 5.7, 0.018, fill=DIV)
        txt(s, desc, cx + 0.2, cy + 1.37, 5.7, 0.55, sz=11, clr=MD)

    fade_tr(s)
    return s

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — FAZIT & CALL TO ACTION (weiß, typografisch stark)
# ─────────────────────────────────────────────────────────────────────────────

def slide_fazit(prs):
    s = new_slide(prs)
    footer(s, 16)

    # Schmale schwarze Linie links — Echo zur Cover-Folie
    rect(s, 0.55, 0.45, 0.04, 6.35, fill=BK)

    t1 = txt(s, "Starten wir\ngemeinsam.", 0.9, 0.5, 11, 2.0, sz=54, clr=BK, fn=FL)
    t2 = txt(s, "ARENA Executive Search ist bereit.", 0.9, 2.5, 10, 0.55,
             sz=22, clr=MD, fn=FL)
    ln = rect(s, 0.9, 3.15, 8.0, 0.025, fill=DIV)

    bullets = [
        "Erfahrener Gründer — 30+ Jahre · 1.500+ Kontakte · 200+ Placements",
        "Klares Geschäftsmodell — Retained Search mit sofort aktivierbarem Deal-Flow",
        "Realistische Zahlen — Break-Even Monat 8, EBIT > €1 Mio. ab Jahr 3",
        "Doppelte Sicherheit — Solo-fähig + Partnerschaft als strategischer Upside",
        "Überschaubarer Kapitalbedarf — € 150.000 mit klar geplanter Tilgung",
    ]
    for i, b in enumerate(bullets):
        txt(s, b, 0.9, 3.35 + i * 0.5, 11.8, 0.42, sz=13, clr=DK, wrap=True)

    rect(s, 0.9, 5.82, 12.0, 0.022, fill=DIV)

    txt(s, "Kontakt", 0.9, 5.97, 2.0, 0.28, sz=10, clr=MD, bold=True)
    txt(s, "Aaron Arena   ·   aaron.arena@arena-executive-search.de   ·   +49 (0) 170 — — — — — —",
        0.9, 6.28, 12.0, 0.3, sz=12, clr=BK)
    txt(s, "www.arena-executive-search.de   ·   LinkedIn: /in/aaronarena",
        0.9, 6.62, 10.0, 0.28, sz=11, clr=MD)

    fade_tr(s)
    apply_anim(s, [
        (t1.shape_id, 'fade', 1,   0, 800),
        (t2.shape_id, 'fade', 1, 700, 600),
        (ln.shape_id, 'fade', 1,1100, 400),
    ])
    return s

# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Erstelle Präsentation …")
    slide_cover(prs)       # 1
    slide_agenda(prs)      # 2
    slide_chance(prs)      # 3
    slide_founder(prs)     # 4
    slide_markt(prs)       # 5
    slide_modell(prs)      # 6
    slide_leistungen(prs)  # 7
    slide_usp(prs)         # 8
    slide_partner(prs)     # 9
    slide_wachstum(prs)    # 10
    slide_finanzen(prs)    # 11
    slide_invest(prs)      # 12
    slide_risiken(prs)     # 13
    slide_roadmap(prs)     # 14
    slide_branchen(prs)    # 15
    slide_fazit(prs)       # 16

    out = "/home/user/ARENA_Executive_Search_Businessplan_2026.pptx"
    prs.save(out)
    size = os.path.getsize(out)
    print(f"✓  Gespeichert: {out}")
    print(f"   Folien: {len(prs.slides)}  |  Größe: {size/1024:.0f} KB")
    return out

if __name__ == "__main__":
    import os
    main()
